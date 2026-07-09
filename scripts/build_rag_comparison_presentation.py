from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))


COLORS = {
    "ink": RGBColor(20, 24, 32),
    "muted": RGBColor(91, 102, 118),
    "bg": RGBColor(247, 249, 252),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(219, 225, 235),
    "blue": RGBColor(18, 111, 181),
    "blue_dark": RGBColor(11, 63, 120),
    "blue_soft": RGBColor(231, 243, 255),
    "teal": RGBColor(18, 140, 131),
    "teal_soft": RGBColor(228, 248, 244),
    "green": RGBColor(35, 151, 88),
    "green_soft": RGBColor(231, 248, 237),
    "red": RGBColor(209, 65, 65),
    "red_soft": RGBColor(255, 237, 237),
    "yellow": RGBColor(255, 214, 71),
    "yellow_soft": RGBColor(255, 247, 219),
    "purple": RGBColor(113, 92, 207),
    "purple_soft": RGBColor(242, 239, 255),
}


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def set_font(run, size: int, *, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str = "",
    size: int = 16,
    *,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, bold=bold, color=color or COLORS["ink"])
    return box


def add_round_rect(slide, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1)
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor, width: int = 1):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.62, 0.36, 11.9, 0.48, title, 28, bold=True, color=COLORS["ink"])
    if subtitle:
        add_textbox(slide, 0.66, 0.88, 11.7, 0.3, subtitle, 12, color=COLORS["muted"])
    add_line(slide, 0.64, 1.24, 12.72, 1.24, COLORS["line"])


def add_footer(slide, number: int) -> None:
    add_textbox(slide, 0.66, 7.14, 7.8, 0.22, "RAG 사용/미사용 추천 검증 비교 · 2026.07.10", 8, color=COLORS["muted"])
    add_textbox(slide, 12.18, 7.14, 0.48, 0.22, f"{number:02d}", 8, bold=True, color=COLORS["blue"], align=PP_ALIGN.RIGHT)


def add_badge(slide, x: float, y: float, text: str, *, fill: RGBColor, color: RGBColor, w: float | None = None):
    width = w or max(1.0, min(3.4, 0.12 * len(text) + 0.5))
    add_round_rect(slide, x, y, width, 0.34, fill, fill)
    add_textbox(slide, x + 0.08, y + 0.055, width - 0.16, 0.18, text, 8, bold=True, color=color, align=PP_ALIGN.CENTER)


def add_metric(slide, x: float, y: float, w: float, h: float, value: str, label: str, *, fill: RGBColor, accent: RGBColor):
    add_round_rect(slide, x, y, w, h, fill, COLORS["line"])
    add_textbox(slide, x + 0.14, y + 0.16, w - 0.28, 0.5, value, 28, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.16, y + 0.78, w - 0.32, 0.3, label, 10, bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def add_bullets(slide, x: float, y: float, w: float, items: list[str], *, size: int = 13, gap: float = 0.45, color: RGBColor | None = None):
    for index, item in enumerate(items):
        yy = y + index * gap
        add_round_rect(slide, x, yy + 0.1, 0.08, 0.08, COLORS["blue"], COLORS["blue"])
        add_textbox(slide, x + 0.22, yy, w - 0.22, 0.28, item, size, color=color or COLORS["ink"])


def format_percent(value: Any) -> str:
    if value is None:
        return "-"
    return f"{float(value) * 100:.1f}%"


def check_text(summary: dict[str, Any], prefix: str) -> str:
    return f"{summary.get(prefix + '_passed_checks')}/{summary.get(prefix + '_total_checks')}"


def case_pass_text(summary: dict[str, Any], prefix: str) -> str:
    return f"{summary.get(prefix + '_passed_cases')}/{summary.get('scenario_cases')}"


def add_result_bar(slide, x: float, y: float, label: str, rate: float, *, accent: RGBColor):
    add_textbox(slide, x, y, 2.2, 0.28, label, 12, bold=True, color=COLORS["ink"])
    add_round_rect(slide, x + 2.4, y + 0.07, 4.65, 0.18, COLORS["line"], COLORS["line"])
    add_round_rect(slide, x + 2.4, y + 0.07, max(0.02, 4.65 * rate), 0.18, accent, accent)
    add_textbox(slide, x + 7.25, y - 0.02, 1.1, 0.25, format_percent(rate), 12, bold=True, color=accent, align=PP_ALIGN.RIGHT)


def add_simple_table(slide, x: float, y: float, w: float, h: float, headers: list[str], rows: list[list[str]], widths: list[float]):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["blue_dark"]
        cell.text = header
        style_cell(cell, 8, COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["bg"] if row_index % 2 else COLORS["white"]
            cell.text = value
            style_cell(cell, 7, COLORS["ink"], align=PP_ALIGN.CENTER if col_index in {1, 2} else PP_ALIGN.LEFT)
    return table_shape


def style_cell(cell, size: int, color: RGBColor, *, bold: bool = False, align: PP_ALIGN | None = None):
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    for paragraph in cell.text_frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            set_font(run, size, bold=bold, color=color)


def compact_failures(failures: list[str], limit: int = 4) -> str:
    names = [item.split(":", 1)[-1] for item in failures[:limit]]
    if len(failures) > limit:
        names.append(f"+{len(failures) - limit}")
    return ", ".join(names)


def build_deck(comparison: dict[str, Any], no_rag: dict[str, Any], *, output_path: Path) -> Path:
    summary = comparison["summary"]
    official = summary["official_course_slots"]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    add_badge(slide, 0.76, 0.76, "추천 품질 검증", fill=COLORS["blue_soft"], color=COLORS["blue_dark"], w=1.55)
    add_textbox(slide, 0.76, 1.34, 7.3, 0.68, "RAG 사용/미사용 비교", 38, bold=True, color=COLORS["ink"])
    add_textbox(slide, 0.79, 2.12, 7.1, 0.42, "제주 접근성 여행 추천의 근거 기반 검증 결과", 19, bold=True, color=COLORS["blue_dark"])
    add_textbox(slide, 0.8, 2.85, 6.65, 0.6, "같은 5개 사용자 상황을 기준으로 RAG 추천과 무RAG 통제 기준선을 동일한 검증표로 비교했습니다.", 15, color=COLORS["muted"])
    add_metric(slide, 0.82, 4.05, 2.25, 1.25, "59/59", "RAG 체크 통과", fill=COLORS["white"], accent=COLORS["green"])
    add_metric(slide, 3.38, 4.05, 2.25, 1.25, "5/59", "무RAG 체크 통과", fill=COLORS["white"], accent=COLORS["red"])
    add_metric(slide, 5.94, 4.05, 2.25, 1.25, "62/62", "공식 코스 슬롯 매칭", fill=COLORS["white"], accent=COLORS["blue"])
    add_round_rect(slide, 8.8, 0.0, 4.53, 7.5, COLORS["blue_dark"], COLORS["blue_dark"])
    add_textbox(slide, 9.15, 1.25, 3.6, 0.5, "핵심 결론", 22, bold=True, color=COLORS["white"])
    add_bullets(
        slide,
        9.2,
        2.1,
        3.5,
        [
            "RAG는 케이스 5/5 통과",
            "무RAG 기준선은 0/5 통과",
            "실패는 도보 부담, 제외 규칙, 편의시설 근거에서 집중",
            "공식 코스 대조는 RAG에서만 가능",
        ],
        size=13,
        gap=0.62,
        color=COLORS["white"],
    )
    add_footer(slide, 1)

    # 2. Experiment design
    slide = prs.slides.add_slide(blank)
    add_title(slide, "비교 설계", "같은 입력, 같은 검증표, 다른 근거 사용 방식")
    columns = [
        (
            "RAG 사용",
            COLORS["green_soft"],
            [
                "장소 카드와 검수 상태 검색",
                "상황별 제외/감점 정책 적용",
                "공식 추천코스 슬롯 대조",
                "추천 근거와 방문 전 확인 항목 생성",
            ],
        ),
        (
            "RAG 미사용",
            COLORS["red_soft"],
            [
                "근거 저장소 없이 일반 제주 대표지 중심 응답",
                "장소별 편의시설·검수 상태 미연결",
                "차단 후보와 감점 사유 추적 불가",
                "통제 기준선으로 재현해 동일 검증표 적용",
            ],
        ),
    ]
    for index, (title, fill, items) in enumerate(columns):
        x = 0.9 + index * 6.0
        add_round_rect(slide, x, 1.72, 5.35, 3.55, fill, COLORS["line"])
        add_textbox(slide, x + 0.28, 2.02, 4.7, 0.42, title, 22, bold=True, color=COLORS["ink"])
        add_bullets(slide, x + 0.35, 2.78, 4.7, items, size=12, gap=0.48)
    add_round_rect(slide, 1.18, 5.82, 10.9, 0.7, COLORS["yellow_soft"], COLORS["line"])
    add_textbox(slide, 1.45, 6.02, 10.35, 0.26, "주의: 무RAG 자료는 실제 외부 모델 운영 로그가 아니라 위험 비교용 통제 기준선입니다.", 14, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # 3. Evidence dataset
    slide = prs.slides.add_slide(blank)
    add_title(slide, "검증 데이터", "상황별 추천 정책과 공식 코스 데이터를 함께 사용")
    metric_cards = [
        ("5", "사용자 상황"),
        ("59", "검증 체크"),
        (str(official.get("courses")), "공식 코스"),
        (str(official.get("stops")), "코스 슬롯"),
        (str(official.get("unmatched_places")), "미매칭 장소"),
    ]
    for index, (value, label) in enumerate(metric_cards):
        add_metric(slide, 0.72 + index * 2.5, 1.62, 2.12, 1.12, value, label, fill=COLORS["white"], accent=COLORS["blue"])
    cases = comparison["cases"]
    rows = [
        [
            case["label"],
            f"{case['with_rag']['passed_checks']}/{case['with_rag']['total_checks']}",
            f"{case['without_rag_baseline']['passed_checks']}/{case['without_rag_baseline']['total_checks']}",
            ", ".join(case["without_rag_baseline"].get("route_categories", [])[:4]),
        ]
        for case in cases
    ]
    add_simple_table(
        slide,
        0.76,
        3.2,
        11.85,
        2.35,
        ["상황", "RAG", "무RAG", "무RAG 추천 유형"],
        rows,
        [1.5, 1.1, 1.15, 8.1],
    )
    add_textbox(slide, 1.0, 6.1, 11.35, 0.34, "평가 기준: 추천 수, 총점, 차단 여부, 도보 부담, 제외 유형, 편의시설, 설명 필수어, 방문 전 확인, 정책 효과", 11, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # 4. Quantitative comparison
    slide = prs.slides.add_slide(blank)
    add_title(slide, "정량 결과", "RAG는 모든 체크를 통과했고, 무RAG는 대부분의 검증 항목에서 실패")
    add_metric(slide, 0.88, 1.55, 2.55, 1.25, case_pass_text(summary, "with_rag"), "RAG 케이스 통과", fill=COLORS["green_soft"], accent=COLORS["green"])
    add_metric(slide, 3.7, 1.55, 2.55, 1.25, check_text(summary, "with_rag"), "RAG 체크 통과", fill=COLORS["green_soft"], accent=COLORS["green"])
    add_metric(slide, 7.08, 1.55, 2.55, 1.25, f"{summary['without_rag_passed_cases']}/{summary['scenario_cases']}", "무RAG 케이스 통과", fill=COLORS["red_soft"], accent=COLORS["red"])
    add_metric(slide, 9.9, 1.55, 2.55, 1.25, check_text(summary, "without_rag"), "무RAG 체크 통과", fill=COLORS["red_soft"], accent=COLORS["red"])
    add_result_bar(slide, 1.1, 3.65, "RAG 케이스 통과율", summary["with_rag_case_pass_rate"], accent=COLORS["green"])
    add_result_bar(slide, 1.1, 4.25, "무RAG 케이스 통과율", summary["without_rag_case_pass_rate"], accent=COLORS["red"])
    add_result_bar(slide, 1.1, 5.0, "RAG 체크 통과율", summary["with_rag_check_pass_rate"], accent=COLORS["green"])
    add_result_bar(slide, 1.1, 5.6, "무RAG 체크 통과율", summary["without_rag_check_pass_rate"], accent=COLORS["red"])
    add_round_rect(slide, 9.72, 3.52, 2.55, 2.42, COLORS["blue_soft"], COLORS["line"])
    add_textbox(slide, 9.98, 3.88, 2.0, 0.42, "차이", 19, bold=True, color=COLORS["blue_dark"], align=PP_ALIGN.CENTER)
    add_textbox(slide, 9.95, 4.58, 2.1, 0.42, "+54개", 30, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)
    add_textbox(slide, 9.95, 5.18, 2.1, 0.32, "체크 통과 격차", 11, bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, 4)

    # 5. Case-by-case comparison
    slide = prs.slides.add_slide(blank)
    add_title(slide, "케이스별 비교", "무RAG는 대표 관광지 추천으로 상황별 제약을 놓친다")
    rows = []
    for case in cases:
        no_rag_case = case["without_rag_baseline"]
        rows.append(
            [
                case["label"],
                f"{case['with_rag']['passed_checks']}/{case['with_rag']['total_checks']}",
                f"{no_rag_case['passed_checks']}/{no_rag_case['total_checks']}",
                ", ".join(no_rag_case.get("route_names", [])[:4]),
                compact_failures(no_rag_case.get("failed_check_names", []), limit=3),
            ]
        )
    add_simple_table(
        slide,
        0.54,
        1.55,
        12.25,
        4.75,
        ["상황", "RAG", "무RAG", "무RAG 추천 경로", "주요 실패"],
        rows,
        [1.15, 0.82, 0.86, 5.65, 3.77],
    )
    add_footer(slide, 5)

    # 6. Failure pattern
    slide = prs.slides.add_slide(blank)
    add_title(slide, "무RAG 실패 패턴", "실패는 한두 장소의 문제가 아니라 근거 부재에서 반복된다")
    patterns = [
        ("점수/차단 검증 불가", "총점 기준과 차단 장소 여부를 재현 가능한 로그로 확인하기 어렵다.", COLORS["red_soft"]),
        ("도보 부담 미반영", "오름·해변·섬처럼 걷기 부담이 큰 장소가 상위 추천에 반복 포함된다.", COLORS["yellow_soft"]),
        ("제외 규칙 위반", "음식 제한, 날씨 민감, 회복기 조건에서 시장·식당·해안·오름이 남는다.", COLORS["red_soft"]),
        ("편의시설 근거 없음", "휠체어 접근, 장애인 화장실, 주차, 휴식 공간 상태를 출처로 대조할 수 없다.", COLORS["purple_soft"]),
        ("방문 전 확인 누락", "경사, 바닥, 주차, 강풍, 그늘, 식사 같은 필수 확인어가 누락된다.", COLORS["blue_soft"]),
    ]
    for index, (title, body, fill) in enumerate(patterns):
        x = 0.75 + (index % 2) * 6.05
        y = 1.62 + (index // 2) * 1.48
        add_round_rect(slide, x, y, 5.52, 1.05, fill, COLORS["line"])
        add_textbox(slide, x + 0.22, y + 0.16, 4.9, 0.26, title, 14, bold=True, color=COLORS["ink"])
        add_textbox(slide, x + 0.22, y + 0.52, 4.95, 0.28, body, 10, color=COLORS["muted"])
    add_round_rect(slide, 3.15, 6.08, 7.05, 0.55, COLORS["red"], COLORS["red"])
    add_textbox(slide, 3.38, 6.24, 6.6, 0.18, "결론: 무RAG는 ‘그럴듯한 제주 코스’는 만들지만, 접근성 추천 검증은 통과하지 못한다.", 11, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    add_footer(slide, 6)

    # 7. What RAG changes
    slide = prs.slides.add_slide(blank)
    add_title(slide, "RAG가 만든 차이", "추천 품질은 모델 문장보다 근거 연결과 검증 가능성에서 갈린다")
    flow = [
        ("검색", "장소 카드·공식 코스·검수 상태"),
        ("정책", "제외/감점·상황별 필수어"),
        ("조립", "추천 경로·근거·방문 전 확인"),
        ("검증", "59개 체크와 케이스별 리포트"),
    ]
    for index, (title, body) in enumerate(flow):
        x = 0.84 + index * 3.05
        add_round_rect(slide, x, 2.02, 2.45, 2.08, COLORS["white"], COLORS["line"])
        add_round_rect(slide, x + 0.18, 2.22, 0.44, 0.44, COLORS["blue"], COLORS["blue"])
        add_textbox(slide, x + 0.18, 2.32, 0.44, 0.16, str(index + 1), 10, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.25, 2.9, 1.9, 0.28, title, 17, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.25, 3.33, 1.92, 0.34, body, 10, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        if index < len(flow) - 1:
            add_line(slide, x + 2.45, 3.06, x + 2.92, 3.06, COLORS["blue"], 2)
    add_round_rect(slide, 1.1, 5.18, 11.15, 0.82, COLORS["green_soft"], COLORS["line"])
    add_textbox(slide, 1.38, 5.42, 10.6, 0.25, "사용자에게 중요한 것은 장소명보다 ‘왜 이 조건에서 가능한지’와 ‘무엇을 확인해야 하는지’다.", 15, bold=True, color=COLORS["green"], align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8. Caveat and next validation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "주의사항과 다음 검증", "현재 비교는 통제 기준선이며, 운영 전 실제 무RAG 로그로 재검증한다")
    add_round_rect(slide, 0.85, 1.62, 5.65, 4.45, COLORS["yellow_soft"], COLORS["line"])
    add_textbox(slide, 1.15, 1.92, 4.9, 0.35, "이번 자료의 해석 범위", 19, bold=True, color=COLORS["ink"])
    add_bullets(
        slide,
        1.2,
        2.58,
        4.85,
        [
            "무RAG는 실제 운영 로그가 아닌 통제 기준선",
            "비교 목적은 위험 패턴과 검증 항목을 드러내는 것",
            "실제 모델 응답을 수집하면 같은 채점기로 재평가 가능",
        ],
        size=12,
        gap=0.5,
    )
    add_round_rect(slide, 6.92, 1.62, 5.65, 4.45, COLORS["blue_soft"], COLORS["line"])
    add_textbox(slide, 7.22, 1.92, 4.9, 0.35, "다음 실험 체크리스트", 19, bold=True, color=COLORS["ink"])
    add_bullets(
        slide,
        7.27,
        2.58,
        4.85,
        [
            "동일 5개 입력으로 실제 무RAG 모델 응답 저장",
            "장소명·편의시설·확인 항목 정답 대조",
            "환각, 제외 규칙 위반, 출처 문구 유무 채점",
            "RAG 프롬프트/검색 결과를 버전 고정",
        ],
        size=12,
        gap=0.47,
    )
    add_footer(slide, 8)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a PPT deck for RAG vs no-RAG comparison.")
    parser.add_argument("--comparison-json", default="data/rag_comparison_report.json")
    parser.add_argument("--no-rag-json", default="data/no_rag_baseline_validation_report.json")
    parser.add_argument("--output", default="docs/rag_comparison_presentation_20260710.pptx")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    comparison = load_json(ROOT / args.comparison_json)
    no_rag = load_json(ROOT / args.no_rag_json)
    output = build_deck(comparison, no_rag, output_path=ROOT / args.output)
    print(f"rag_comparison_presentation={output}")
    print(f"slides=8")
    print(
        "summary="
        f"with_rag:{comparison['summary']['with_rag_passed_checks']}/{comparison['summary']['with_rag_total_checks']}, "
        f"without_rag:{comparison['summary']['without_rag_passed_checks']}/{comparison['summary']['without_rag_total_checks']}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
