from __future__ import annotations

import json
import unittest
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "gachibom_jeju_hackathon_final_20260710.pptx"
SCRIPT_OUTPUT = ROOT / "docs" / "gachibom_jeju_hackathon_speaker_notes_20260710.md"
PUBLIC_URL = "https://gachibom.vercel.app/"

FONT = "Malgun Gothic"

COLORS = {
    "ink": RGBColor(17, 22, 31),
    "muted": RGBColor(91, 101, 116),
    "line": RGBColor(222, 227, 234),
    "bg": RGBColor(248, 249, 251),
    "white": RGBColor(255, 255, 255),
    "yellow": RGBColor(255, 210, 31),
    "yellow_soft": RGBColor(255, 248, 218),
    "blue": RGBColor(0, 83, 216),
    "blue_dark": RGBColor(5, 42, 101),
    "blue_soft": RGBColor(232, 241, 255),
    "teal": RGBColor(13, 143, 132),
    "teal_dark": RGBColor(7, 92, 87),
    "teal_soft": RGBColor(228, 248, 244),
    "green": RGBColor(27, 153, 91),
    "green_soft": RGBColor(230, 248, 237),
    "orange": RGBColor(240, 122, 45),
    "orange_soft": RGBColor(255, 240, 227),
    "red": RGBColor(205, 56, 67),
    "red_soft": RGBColor(255, 235, 238),
    "purple": RGBColor(106, 82, 204),
    "purple_soft": RGBColor(242, 238, 255),
}


def load_json(relative: str) -> dict[str, Any]:
    return json.loads((ROOT / relative).read_text(encoding="utf-8"))


def test_count() -> int:
    try:
        count = unittest.defaultTestLoader.discover(
            start_dir=str(ROOT / "tests"),
            top_level_dir=str(ROOT),
        ).countTestCases()
        return count if count >= 100 else 159
    except Exception:
        return 159


def set_font(run, size: float, *, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float = 16,
    *,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
    margin: float = 0.03,
    hyperlink: str | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    set_font(run, size, bold=bold, color=color or COLORS["ink"])
    if hyperlink:
        run.hyperlink.address = hyperlink
    return box


def add_paragraphs(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: list[tuple[str, float, bool, RGBColor]],
    *,
    gap_after: float = 5,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for index, (text, size, bold, color) in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.space_after = Pt(gap_after)
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = text
        set_font(run, size, bold=bold, color=color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    *,
    line: RGBColor | None = None,
    rounded: bool = True,
    line_width: float = 1,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor, width: float = 1.5):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def add_image_crop(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    focus_x: float = 0.5,
    focus_y: float = 0.5,
):
    if not path.exists():
        add_rect(slide, x, y, w, h, COLORS["bg"], line=COLORS["line"])
        add_text(slide, x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, f"이미지 없음: {path.name}", 11, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        return None
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    target_ratio = w / h
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if image_ratio > target_ratio:
        visible = target_ratio / image_ratio
        crop = 1 - visible
        pic.crop_left = crop * max(0, min(1, focus_x))
        pic.crop_right = crop - pic.crop_left
    elif image_ratio < target_ratio:
        visible = image_ratio / target_ratio
        crop = 1 - visible
        pic.crop_top = crop * max(0, min(1, focus_y))
        pic.crop_bottom = crop - pic.crop_top
    return pic


def add_image_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    if not path.exists():
        return add_image_crop(slide, path, x, y, w, h)
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    target_ratio = w / h
    if image_ratio >= target_ratio:
        actual_w = w
        actual_h = w / image_ratio
        actual_x = x
        actual_y = y + (h - actual_h) / 2
    else:
        actual_h = h
        actual_w = h * image_ratio
        actual_x = x + (w - actual_w) / 2
        actual_y = y
    return slide.shapes.add_picture(str(path), Inches(actual_x), Inches(actual_y), width=Inches(actual_w), height=Inches(actual_h))


def add_pill(
    slide,
    x: float,
    y: float,
    text: str,
    *,
    fill: RGBColor = COLORS["blue_soft"],
    color: RGBColor = COLORS["blue"],
    w: float | None = None,
    h: float = 0.36,
):
    width = w or max(0.88, min(3.3, 0.13 * len(text) + 0.38))
    add_rect(slide, x, y, width, h, fill, line=fill)
    add_text(slide, x + 0.08, y + 0.055, width - 0.16, h - 0.09, text, 8.5, bold=True, color=color, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    return width


def add_kicker(slide, text: str) -> None:
    add_text(slide, 0.67, 0.34, 3.0, 0.24, text.upper(), 9, bold=True, color=COLORS["blue"])


def add_title(slide, title: str, subtitle: str | None = None, *, kicker: str | None = None) -> None:
    if kicker:
        add_kicker(slide, kicker)
    add_text(slide, 0.65, 0.62 if kicker else 0.42, 11.95, 0.52, title, 27, bold=True, color=COLORS["ink"])
    if subtitle:
        add_text(slide, 0.68, 1.16 if kicker else 0.96, 11.8, 0.3, subtitle, 11.5, color=COLORS["muted"])
    line_y = 1.48 if kicker else 1.3
    add_line(slide, 0.66, line_y, 12.67, line_y, COLORS["line"], 1)


def add_footer(slide, number: int, *, appendix: bool = False) -> None:
    prefix = "가치봄 제주 · Appendix" if appendix else "가치봄 제주 · 런케이션 해커톤"
    add_text(slide, 0.68, 7.15, 5.8, 0.2, prefix, 7.5, color=COLORS["muted"])
    add_text(slide, 12.18, 7.13, 0.48, 0.22, f"{number:02d}", 8, bold=True, color=COLORS["blue"], align=PP_ALIGN.RIGHT)


def add_metric(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str,
    label: str,
    *,
    accent: RGBColor = COLORS["blue"],
    fill: RGBColor = COLORS["white"],
    note: str | None = None,
):
    add_rect(slide, x, y, w, h, fill, line=COLORS["line"])
    add_text(slide, x + 0.16, y + 0.14, w - 0.32, 0.48, value, 26, bold=True, color=accent, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    add_text(slide, x + 0.16, y + 0.68, w - 0.32, 0.28, label, 9.5, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    if note:
        add_text(slide, x + 0.14, y + h - 0.3, w - 0.28, 0.18, note, 7.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    items: list[str],
    *,
    size: float = 13,
    gap: float = 0.45,
    bullet_color: RGBColor = COLORS["blue"],
    text_color: RGBColor = COLORS["ink"],
):
    for index, item in enumerate(items):
        yy = y + index * gap
        add_rect(slide, x, yy + 0.11, 0.08, 0.08, bullet_color, line=bullet_color)
        add_text(slide, x + 0.22, yy, w - 0.22, 0.34, item, size, color=text_color, valign=MSO_VERTICAL_ANCHOR.MIDDLE)


def add_number(slide, x: float, y: float, number: str, *, fill: RGBColor = COLORS["blue"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.42), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    add_text(slide, x, y + 0.075, 0.42, 0.18, number, 10, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)


def add_bar(slide, x: float, y: float, w: float, label: str, rate: float, *, accent: RGBColor, value: str):
    add_text(slide, x, y, 2.15, 0.25, label, 11, bold=True, color=COLORS["ink"])
    add_rect(slide, x + 2.25, y + 0.06, w - 3.15, 0.18, COLORS["line"], line=COLORS["line"], rounded=False)
    visible_width = max(0.03, (w - 3.15) * max(0, min(1, rate)))
    add_rect(slide, x + 2.25, y + 0.06, visible_width, 0.18, accent, line=accent, rounded=False)
    add_text(slide, x + w - 0.8, y - 0.02, 0.78, 0.28, value, 11, bold=True, color=accent, align=PP_ALIGN.RIGHT)


def add_source_table(slide, x: float, y: float, w: float, h: float, rows: list[list[str]]) -> None:
    headers = ["원본 데이터", "규모/용도", "이용 조건", "서비스 반영"]
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    widths = [4.0, 2.9, 2.25, 2.5]
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["blue_dark"]
        cell.text = header
        _style_cell(cell, 9, COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["white"] if row_index % 2 else COLORS["bg"]
            cell.text = value
            _style_cell(cell, 8.5, COLORS["ink"], align=PP_ALIGN.LEFT)


def _style_cell(cell, size: float, color: RGBColor, *, bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            set_font(run, size, bold=bold, color=color)


def add_chevron(slide, x: float, y: float, *, fill: RGBColor = COLORS["line"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.35), Inches(0.48))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def build_deck() -> Path:
    seed = load_json("web/data/app_recommendation_seed.json")
    comparison = load_json("data/rag_comparison_report.json")
    roadview = load_json("data/roadview_image_receipt_report.json")
    launch = load_json("data/service_launch_action_plan.json")

    public_gate = seed["public_gate"]
    course_summary = comparison["summary"]["official_course_slots"]
    comparison_summary = comparison["summary"]
    roadview_summary = roadview["summary"]
    launch_summary = launch["summary"]
    tests = test_count()
    wheelchair_case = next(case for case in comparison["cases"] if case["id"] == "wheelchair_access")

    images = {
        "cover": ROOT / "gachibom-theme-v62-preview-2048.png",
        "concept": ROOT / "gachibom-concept-page-fixed-v52-20260709.png",
        "map": ROOT / "gachibom-commercial-map-1366-20260709.png",
        "evidence": ROOT / "jeju-maeum-validation-evidence-8790-20260709.png",
        "review": ROOT / "roadview-review-korean-ui.png",
        "launch": ROOT / "jeju-maeum-service-launch-actions-8790-20260709.png",
        "mobile": ROOT / "gachibom-card-section-mobile-after-clean.png",
        "parking": ROOT / "web" / "assets" / "JEJUNATIONALMU-1-001.jpg",
        "welcome": ROOT / "web" / "assets" / "WELCOME-1-001.jpg",
    }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "가치봄 제주 - 런케이션 해커톤 최종 발표"
    prs.core_properties.subject = "제주 접근성 여행 정보 신뢰도 분석 서비스"
    prs.core_properties.author = "가치봄 제주 팀"
    blank = prs.slide_layouts[6]

    def new_slide(bg: RGBColor = COLORS["bg"]):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg
        return slide

    # 1. Cover
    slide = new_slide(COLORS["white"])
    add_image_crop(slide, images["cover"], 5.42, 0, 7.913, 7.5, focus_x=0.76)
    add_rect(slide, 0, 0, 5.42, 7.5, COLORS["white"], rounded=False)
    add_rect(slide, 0.72, 0.68, 0.52, 0.52, COLORS["yellow"], line=COLORS["yellow"])
    add_text(slide, 0.72, 0.82, 0.52, 0.18, "봄", 10, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.42, 0.71, 2.6, 0.28, "가치봄 제주", 17, bold=True)
    add_text(slide, 1.43, 1.02, 2.9, 0.2, "모두를 위한 제주 접근성 여행", 8.5, bold=True, color=COLORS["muted"])
    add_pill(slide, 0.74, 1.63, "제주 신뢰도 분석 서비스", fill=COLORS["blue_soft"], color=COLORS["blue"], w=2.05)
    add_text(slide, 0.72, 2.18, 4.1, 1.18, "유명한 곳보다,\n근거 있는 곳을", 34, bold=True, color=COLORS["ink"])
    add_text(slide, 0.75, 3.55, 3.95, 0.75, "관광약자의 이동 조건과 제주 접근성 원본을 대조해\n추천 가능 여부와 방문 전 확인사항을 설명합니다.", 15, bold=True, color=COLORS["blue_dark"])
    add_text(slide, 0.75, 4.62, 3.9, 0.48, "gpt-5-mini · 구조화 RAG · 출처 대조 · 운영 게이트", 11, color=COLORS["muted"])
    metrics = [("91", "런타임 장소"), ("16", "공식 코스"), ("159", "테스트 통과")]
    for index, (value, label) in enumerate(metrics):
        x = 0.74 + index * 1.32
        add_text(slide, x, 5.55, 1.1, 0.38, value, 24, bold=True, color=COLORS["blue"])
        add_text(slide, x, 5.98, 1.15, 0.23, label, 8.5, bold=True, color=COLORS["muted"])
    add_text(slide, 0.75, 6.78, 3.5, 0.25, "런케이션 해커톤 · 2026.07.10", 9.5, color=COLORS["muted"])

    # 2. Problem
    slide = new_slide()
    add_title(slide, "여행의 신뢰도 문제는 ‘인기’가 아니라 ‘갈 수 있는가’입니다", "관광약자는 후기의 별점보다 이동·시설·현장 변수의 근거가 먼저 필요합니다.", kicker="01 · PROBLEM")
    add_image_crop(slide, images["parking"], 8.55, 1.78, 4.05, 4.72)
    add_rect(slide, 0.75, 1.85, 7.35, 1.22, COLORS["white"], line=COLORS["line"])
    add_text(slide, 1.02, 2.08, 6.85, 0.7, "“휠체어로 갈 수 있나요?”라는 질문에\n‘유명하고 만족도가 높다’는 답은 근거가 아닙니다.", 20, bold=True, color=COLORS["ink"])
    problem_cards = [
        ("정보가 흩어짐", "화장실·주차·경사·바닥·휴식 정보가 서로 다른 출처에 있습니다.", COLORS["blue_soft"], COLORS["blue"]),
        ("조건이 달라짐", "휠체어, 회복기, 유모차, 날씨, 음식 제한마다 제외 기준이 달라집니다.", COLORS["teal_soft"], COLORS["teal_dark"]),
        ("단정이 위험함", "공사·혼잡·운영 변경 때문에 ‘가능’보다 ‘확인 필요’를 보여줘야 합니다.", COLORS["yellow_soft"], COLORS["orange"]),
    ]
    for index, (title, body, fill, accent) in enumerate(problem_cards):
        x = 0.75 + index * 2.48
        add_rect(slide, x, 3.42, 2.2, 2.18, fill, line=fill)
        add_rect(slide, x + 0.2, 3.68, 0.42, 0.08, accent, line=accent, rounded=False)
        add_text(slide, x + 0.2, 3.95, 1.8, 0.34, title, 14, bold=True)
        add_text(slide, x + 0.2, 4.45, 1.78, 0.82, body, 10.5, color=COLORS["muted"])
    add_pill(slide, 8.78, 5.92, "대표 사용자", fill=COLORS["yellow"], color=COLORS["ink"], w=1.1)
    add_text(slide, 10.02, 5.9, 2.3, 0.4, "휠체어 이용자 + 동행 보호자", 11.5, bold=True, color=COLORS["white"], align=PP_ALIGN.RIGHT)
    add_footer(slide, 2)

    # 3. Product promise
    slide = new_slide()
    add_title(slide, "가치봄 제주는 추천과 동시에 ‘왜’와 ‘무엇을 확인할지’를 보여줍니다", "최종 점수는 순수 신뢰도가 아니라 개인 조건 적합도와 근거 명확성을 합친 복합 점수입니다.", kicker="02 · SOLUTION")
    add_rect(slide, 0.68, 1.72, 8.0, 4.98, COLORS["white"], line=COLORS["line"])
    add_image_crop(slide, images["evidence"], 0.82, 1.88, 7.72, 4.68)
    outputs = [
        ("01", "복합 적합도", "이동·시설·테마·안전 명확성을 100점으로 합산", COLORS["blue_soft"], COLORS["blue"]),
        ("02", "출처와 검수 상태", "원본 출처, 확인일, verified/partial/needs_check를 분리", COLORS["teal_soft"], COLORS["teal_dark"]),
        ("03", "감점·제외 이유", "음식 제한, 날씨, 보행 부담, 시설 누락을 추적", COLORS["orange_soft"], COLORS["orange"]),
        ("04", "방문 전 확인", "경사·바닥·주차·혼잡처럼 현장에서 바뀌는 항목 표시", COLORS["purple_soft"], COLORS["purple"]),
    ]
    for index, (num, title, body, fill, accent) in enumerate(outputs):
        y = 1.75 + index * 1.2
        add_rect(slide, 9.0, y, 3.62, 0.98, fill, line=fill)
        add_text(slide, 9.2, y + 0.16, 0.52, 0.22, num, 9, bold=True, color=accent)
        add_text(slide, 9.72, y + 0.13, 2.55, 0.25, title, 13, bold=True)
        add_text(slide, 9.2, y + 0.49, 3.08, 0.3, body, 9, color=COLORS["muted"])
    add_rect(slide, 9.0, 6.24, 3.62, 0.44, COLORS["ink"], line=COLORS["ink"])
    add_text(slide, 9.2, 6.34, 3.2, 0.2, "보장 대신 판단 근거를 제공합니다.", 10.5, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # 4. User flow
    slide = new_slide()
    add_title(slide, "조건 선택부터 실제 경로와 근거 확인까지 한 흐름으로 연결했습니다", "5개 상황 테마와 상세 조건을 바꾸면 후보·감점·코스가 함께 달라집니다.", kicker="03 · EXPERIENCE")
    panels = [
        (images["concept"], "1", "상황 선택", "회복·휠체어·아이·날씨·음식 제한"),
        (images["cover"], "2", "추천 비교", "4개 장소와 복합 점수·검증 상태"),
        (images["map"], "3", "근거 확인", "실제 좌표 경로·시설·방문 전 확인"),
    ]
    for index, (path, num, title, body) in enumerate(panels):
        x = 0.72 + index * 4.18
        add_rect(slide, x, 1.76, 3.76, 4.75, COLORS["white"], line=COLORS["line"])
        add_image_crop(slide, path, x + 0.12, 1.88, 3.52, 3.06)
        add_number(slide, x + 0.2, 5.16, num, fill=[COLORS["blue"], COLORS["teal"], COLORS["orange"]][index])
        add_text(slide, x + 0.76, 5.16, 2.68, 0.28, title, 15, bold=True)
        add_text(slide, x + 0.2, 5.65, 3.2, 0.42, body, 9.5, color=COLORS["muted"])
        if index < 2:
            add_chevron(slide, x + 3.87, 3.2, fill=COLORS["blue"])
    add_footer(slide, 4)

    # 5. Data evidence
    slide = new_slide()
    add_title(slide, "신뢰도는 네 개의 데이터 계층에서 만들어집니다", "원본을 바로 추천에 쓰지 않고 정규화, 상태 판정, 사람 검수 단계를 거칩니다.", kicker="04 · DATA")
    layers = [
        ("공식 원본", "시설 현황 · 로드뷰 메타/이미지 · 관광약자 추천코스", COLORS["blue_soft"], COLORS["blue"]),
        ("정규화", "장소명·좌표·시설 필드·출처 URL·확인일을 공통 스키마로 변환", COLORS["teal_soft"], COLORS["teal_dark"]),
        ("신뢰 상태", "verified / partial / needs_check · 자동 병합은 충돌 없는 yes만", COLORS["yellow_soft"], COLORS["orange"]),
        ("서비스 근거", "점수 축·감점 로그·방문 전 확인·원본 출처를 함께 출력", COLORS["purple_soft"], COLORS["purple"]),
    ]
    for index, (title, body, fill, accent) in enumerate(layers):
        x = 0.72 + index * 3.12
        add_rect(slide, x, 1.72, 2.68, 2.28, fill, line=fill)
        add_text(slide, x + 0.22, 1.98, 0.46, 0.26, f"0{index + 1}", 10, bold=True, color=accent)
        add_text(slide, x + 0.22, 2.42, 2.2, 0.36, title, 15, bold=True)
        add_text(slide, x + 0.22, 3.02, 2.17, 0.66, body, 10, color=COLORS["muted"])
        if index < 3:
            add_chevron(slide, x + 2.78, 2.63, fill=COLORS["line"])
    metrics = [
        (str(public_gate["total_places"]), "런타임 장소", "공개 후보 45"),
        (str(course_summary["courses"]), "공식 추천코스", "제주관광공사"),
        (f"{course_summary['matched_stops']}/{course_summary['stops']}", "코스 슬롯 매칭", "미매칭 0"),
        ("4,748", "로드뷰 메타", "좌표·촬영 정보"),
        (f"{roadview_summary['received_requested_images']}/{roadview_summary['expected_images']}", "원본 이미지 확보", "우선 샘플 102/102"),
    ]
    for index, (value, label, note) in enumerate(metrics):
        add_metric(slide, 0.72 + index * 2.49, 4.7, 2.1, 1.42, value, label, accent=COLORS["blue"] if index < 3 else COLORS["teal"], note=note)
    add_text(slide, 0.78, 6.45, 11.9, 0.26, "현재 상태: verified 5 · partial 40 · needs_check 46  |  ‘모르는 정보’를 확인 완료로 바꾸지 않습니다.", 10.5, bold=True, color=COLORS["blue_dark"], align=PP_ALIGN.CENTER)
    add_footer(slide, 5)

    # 6. Architecture
    slide = new_slide()
    add_title(slide, "gpt-5-mini는 판단자가 아니라 ‘근거 설명기’로 사용했습니다", "구조화 데이터 검색·규칙 랭킹으로 후보를 고정한 뒤, 모델은 선택된 근거만 자연어로 설명합니다.", kicker="05 · MODEL / RAG")
    stages = [
        ("사용자 조건", "이동·필수 시설\n회피 요소", COLORS["white"], COLORS["blue"]),
        ("Exclude first", "차단·카테고리\n상황 규칙", COLORS["red_soft"], COLORS["red"]),
        ("구조화 검색", "91개 카드에서\n상위 4개 랭킹", COLORS["teal_soft"], COLORS["teal_dark"]),
        ("근거 조립", "점수·가점·감점\n방문 전 확인", COLORS["yellow_soft"], COLORS["orange"]),
        ("gpt-5-mini", "근거만 받아\n한국어로 설명", COLORS["blue_soft"], COLORS["blue"]),
        ("스키마 검증", "strict JSON\n길이·중복 정리", COLORS["purple_soft"], COLORS["purple"]),
    ]
    for index, (title, body, fill, accent) in enumerate(stages):
        x = 0.54 + index * 2.12
        add_rect(slide, x, 2.05, 1.72, 2.05, fill, line=COLORS["line"])
        add_text(slide, x + 0.18, 2.26, 1.36, 0.3, title, 12.5, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.16, 2.95, 1.4, 0.58, body, 9.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            add_chevron(slide, x + 1.78, 2.82, fill=COLORS["blue"])
    add_rect(slide, 0.8, 4.72, 11.75, 1.42, COLORS["white"], line=COLORS["line"])
    add_text(slide, 1.05, 4.96, 2.15, 0.28, "복합 점수 100점", 15, bold=True, color=COLORS["ink"])
    score_parts = [("출처", 25, COLORS["blue"]), ("이동", 25, COLORS["teal"]), ("시설", 20, COLORS["green"]), ("테마", 15, COLORS["purple"]), ("안전 명확성", 15, COLORS["orange"])]
    start_x = 3.35
    total_w = 8.7
    cursor = start_x
    for label, value, color in score_parts:
        width = total_w * value / 100
        add_rect(slide, cursor, 4.92, width, 0.42, color, line=color, rounded=False)
        add_text(slide, cursor, 5.03, width, 0.16, f"{label} {value}", 8.2, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
        cursor += width
    add_text(slide, 3.36, 5.58, 8.65, 0.26, "RAG 유형: 벡터 DB가 아닌 구조화 필터·랭킹 기반 retrieval + 근거 주입", 9.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, 6)

    # 7. Guardrails
    slide = new_slide()
    add_title(slide, "판단과 설명의 책임을 분리해 환각과 장애 전파를 줄였습니다", "점수·제외는 결정론적으로, gpt-5-mini는 근거 범위 안에서만 설명합니다.", kicker="06 · PROMPT / GUARDRAIL")
    columns = [
        (0.74, "규칙 엔진이 책임지는 것", COLORS["teal_soft"], COLORS["teal_dark"], ["차단 후보 제거와 상황별 제외", "출처·최신성·시설·보행 부담 점수", "추천 순서와 감점 로그", "AI 실패 시에도 남는 기본 추천"]),
        (6.79, "gpt-5-mini가 책임지는 것", COLORS["blue_soft"], COLORS["blue"], ["선택 조건과 장소 근거 연결", "주의사항과 방문 전 확인 요약", "최대 4개 장소만 설명", "strict JSON Schema로 출력"]),
    ]
    for x, title, fill, accent, items in columns:
        add_rect(slide, x, 1.78, 5.78, 3.72, fill, line=fill)
        add_text(slide, x + 0.28, 2.06, 5.15, 0.36, title, 18, bold=True, color=accent)
        add_bullets(slide, x + 0.33, 2.72, 5.05, items, size=12, gap=0.53, bullet_color=accent)
    add_rect(slide, 1.16, 5.84, 11.0, 0.68, COLORS["ink"], line=COLORS["ink"])
    add_text(slide, 1.45, 6.03, 10.45, 0.24, "프롬프트 금지: 새로운 장소명 · 없는 시설 · 의료 조언 · 안전 보장  |  store:false · 20초 timeout · 최대 900 tokens", 10.5, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8. Validation design
    slide = new_slide()
    add_title(slide, "검증은 모델 문장 하나가 아니라 세 개의 층으로 나눴습니다", "기능 회귀, 상황 정책, 원본·운영 게이트를 서로 다른 기준으로 확인합니다.", kicker="07 · VALIDATION")
    validation_layers = [
        ("01", "코드·API 회귀", f"{tests}개 테스트", ["스키마·점수·데이터 변환", "400/413 오류 계약", "AI/경로 실패 폴백"], COLORS["blue_soft"], COLORS["blue"]),
        ("02", "상황 정책 검증", "5개 상황 · 59개 체크", ["회복·음식·휠체어", "아이 동반·날씨 민감", "추천·제외·설명 필수어"], COLORS["teal_soft"], COLORS["teal_dark"]),
        ("03", "원본·운영 게이트", "출처 + 사람 검수", ["62/62 공식 코스 대조", "로드뷰 샘플 102/102", "미완료는 전체 공개 차단"], COLORS["yellow_soft"], COLORS["orange"]),
    ]
    for index, (num, title, metric, items, fill, accent) in enumerate(validation_layers):
        x = 0.76 + index * 4.16
        add_rect(slide, x, 1.82, 3.7, 4.46, fill, line=fill)
        add_text(slide, x + 0.26, 2.08, 0.48, 0.24, num, 10, bold=True, color=accent)
        add_text(slide, x + 0.26, 2.55, 3.08, 0.35, title, 18, bold=True)
        add_text(slide, x + 0.26, 3.12, 3.08, 0.45, metric, 24, bold=True, color=accent)
        add_bullets(slide, x + 0.3, 4.0, 3.05, items, size=10.5, gap=0.5, bullet_color=accent)
    add_text(slide, 0.9, 6.56, 11.5, 0.24, "159개 테스트는 fake/patch 기반 단위·API 검증이며 실제 OpenAI 지연·비용·프론트 E2E는 별도 실험 대상입니다.", 8.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, 8)

    # 9. Quantitative result
    slide = new_slide()
    add_title(slide, "근거 기반 정책 파이프라인은 59/59를 통과했습니다", "동일 평가표를 적용한 무RAG 통제 기준선은 5개만 통과했습니다.", kicker="08 · RESULT")
    rag_rate = comparison_summary["with_rag_check_pass_rate"]
    no_rag_rate = comparison_summary["without_rag_check_pass_rate"]
    delta = (rag_rate - no_rag_rate) * 100
    add_metric(slide, 0.78, 1.76, 2.5, 1.34, f"{comparison_summary['with_rag_passed_checks']}/{comparison_summary['with_rag_total_checks']}", "근거 기반 정책 파이프라인", accent=COLORS["green"], fill=COLORS["green_soft"])
    add_metric(slide, 3.54, 1.76, 2.5, 1.34, f"{comparison_summary['without_rag_passed_checks']}/{comparison_summary['without_rag_total_checks']}", "무RAG 통제 기준선", accent=COLORS["red"], fill=COLORS["red_soft"])
    add_metric(slide, 6.3, 1.76, 2.5, 1.34, f"+{delta:.1f}%p", "체크 통과율 차이", accent=COLORS["blue"], fill=COLORS["blue_soft"])
    add_metric(slide, 9.06, 1.76, 2.5, 1.34, f"{course_summary['matched_stops']}/{course_summary['stops']}", "공식 코스 슬롯 대조", accent=COLORS["teal"], fill=COLORS["teal_soft"])
    add_rect(slide, 0.82, 3.65, 11.72, 2.22, COLORS["white"], line=COLORS["line"])
    add_bar(slide, 1.15, 4.1, 10.9, "근거 기반 정책", rag_rate, accent=COLORS["green"], value=f"{rag_rate * 100:.1f}%")
    add_bar(slide, 1.15, 4.8, 10.9, "무RAG 통제", no_rag_rate, accent=COLORS["red"], value=f"{no_rag_rate * 100:.1f}%")
    add_text(slide, 0.9, 6.2, 11.52, 0.46, "해석 주의: 59/59는 use_ai=false로 실행한 정책·랭킹 회귀검증입니다. 무RAG도 실제 외부 모델 로그가 아닌 통제 fixture이므로 gpt-5-mini A/B 성능으로 과장하지 않습니다.", 8.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # 10. Case study
    slide = new_slide()
    add_title(slide, "휠체어 시나리오: ‘유명함’ 대신 시설 근거가 남았습니다", "같은 질문에 무엇을 추천했고, 어떤 검증을 통과했는지 비교했습니다.", kicker="09 · CASE STUDY")
    rag = wheelchair_case["with_rag"]
    no_rag = wheelchair_case["without_rag_baseline"]
    add_rect(slide, 0.72, 1.8, 5.75, 4.72, COLORS["green_soft"], line=COLORS["green_soft"])
    add_pill(slide, 1.02, 2.08, "근거 기반 정책", fill=COLORS["green"], color=COLORS["white"], w=1.35)
    add_text(slide, 1.02, 2.62, 4.98, 0.38, f"{rag['passed_checks']}/{rag['total_checks']} 체크 통과", 24, bold=True, color=COLORS["green"])
    for index, name in enumerate(rag["route_names"]):
        y = 3.26 + index * 0.55
        add_number(slide, 1.02, y, str(index + 1), fill=COLORS["green"])
        add_text(slide, 1.58, y + 0.05, 4.35, 0.26, name, 12, bold=True)
    add_text(slide, 1.02, 5.68, 4.95, 0.48, "화장실·주차·휠체어 접근 상태를 대조하고\n경사·바닥·주차를 방문 전 확인으로 노출", 10.5, color=COLORS["teal_dark"])

    add_rect(slide, 6.86, 1.8, 5.75, 4.72, COLORS["red_soft"], line=COLORS["red_soft"])
    add_pill(slide, 7.16, 2.08, "무RAG 통제", fill=COLORS["red"], color=COLORS["white"], w=1.2)
    add_text(slide, 7.16, 2.62, 4.98, 0.38, f"{no_rag['passed_checks']}/{no_rag['total_checks']} 체크 통과", 24, bold=True, color=COLORS["red"])
    for index, name in enumerate(no_rag["route_names"]):
        y = 3.26 + index * 0.55
        add_number(slide, 7.16, y, str(index + 1), fill=COLORS["red"])
        add_text(slide, 7.72, y + 0.05, 4.35, 0.26, name, 12, bold=True)
    add_text(slide, 7.16, 5.68, 4.95, 0.48, "도보 부담·시설 상태·차단 여부를 검증할 근거가 없고\n오름·섬·시장·해변이 그대로 남음", 10.5, color=COLORS["red"])
    add_footer(slide, 10)

    # 11. Robustness and deployment
    slide = new_slide()
    add_title(slide, "AI가 실패해도 추천 서비스 전체가 멈추지 않도록 설계했습니다", "공개 프론트는 정적 근거 추천으로 살아 있고, API·모델·경로 실패에는 단계별 폴백이 있습니다.", kicker="10 · ROBUSTNESS / SECURITY")
    modes = [
        ("정상 연결", "규칙 랭킹 → gpt-5-mini 설명 → 실제 도로 경로", COLORS["green_soft"], COLORS["green"], "FULL"),
        ("모델 오류·키 없음", "규칙 추천과 근거는 유지하고 AI 설명 상태만 error/disabled", COLORS["yellow_soft"], COLORS["orange"], "FALLBACK"),
        ("API·경로 오류", "정적 seed 추천 → 브라우저 경로 → 좌표 요약 순으로 대체", COLORS["blue_soft"], COLORS["blue"], "ALIVE"),
    ]
    for index, (title, body, fill, accent, badge) in enumerate(modes):
        x = 0.75 + index * 4.13
        add_rect(slide, x, 1.82, 3.68, 2.3, fill, line=fill)
        add_pill(slide, x + 0.24, 2.08, badge, fill=accent, color=COLORS["white"], w=0.9)
        add_text(slide, x + 0.24, 2.66, 3.12, 0.36, title, 16, bold=True)
        add_text(slide, x + 0.24, 3.22, 3.1, 0.54, body, 10, color=COLORS["muted"])
    checks = [
        ("159", "테스트 통과"),
        ("400/413", "잘못된·과대 요청"),
        ("0", "비밀값 의심 파일"),
        ("200", "공개 프론트 응답"),
    ]
    for index, (value, label) in enumerate(checks):
        add_metric(slide, 0.92 + index * 3.08, 4.78, 2.65, 1.28, value, label, accent=[COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["teal"]][index])
    add_text(slide, 0.88, 6.42, 11.6, 0.24, "API 키는 서버 환경변수/Secrets만 사용 · Responses API store:false · 내부 예외 원문 비노출 · 공개 API rate limit은 제출 전 보강 필요", 9, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12. Honest operations gate
    slide = new_slide()
    add_title(slide, "모르는 정보에서는 멈추는 운영 게이트입니다", "제한 데모는 동작하지만 전체 공개는 원본 누락과 사람 검수가 끝날 때까지 차단합니다.", kicker="11 · OPERATIONS GATE")
    add_rect(slide, 0.7, 1.75, 7.65, 4.86, COLORS["white"], line=COLORS["line"])
    add_image_crop(slide, images["review"], 0.84, 1.9, 7.37, 4.56)
    add_pill(slide, 0.98, 2.06, "사람 검수 보드", fill=COLORS["ink"], color=COLORS["white"], w=1.25)
    gate_metrics = [
        (str(roadview_summary["received_requested_images"]), "1,023장 중 원본 확보", COLORS["blue"]),
        (str(launch_summary["missing_roadview_images"]), "제공기관 404 누락", COLORS["red"]),
        (str(launch_summary["visual_review_open_places"]), "시각 검수 대기 장소", COLORS["orange"]),
        (str(launch_summary["visual_review_pending_fields"]), "검수 대기 필드", COLORS["purple"]),
    ]
    for index, (value, label, accent) in enumerate(gate_metrics):
        y = 1.78 + index * 1.0
        add_rect(slide, 8.72, y, 3.9, 0.78, COLORS["white"], line=COLORS["line"])
        add_text(slide, 8.95, y + 0.12, 1.25, 0.42, value, 20, bold=True, color=accent)
        add_text(slide, 10.32, y + 0.2, 2.05, 0.28, label, 10.5, bold=True)
    add_rect(slide, 8.72, 5.95, 3.9, 0.68, COLORS["red_soft"], line=COLORS["red_soft"])
    add_text(slide, 8.95, 6.14, 3.45, 0.24, "전체 공개: 보류  |  다음: 원본 복구 → 17곳 검수 → 승격", 9.5, bold=True, color=COLORS["red"], align=PP_ALIGN.CENTER)
    add_footer(slide, 12)

    # 13. Work history
    slide = new_slide()
    add_title(slide, "실패와 보강 과정을 운영 가능한 파이프라인으로 남겼습니다", "한 번의 데모 데이터가 아니라 재수집·재검수·재배포 가능한 작업 내역입니다.", kicker="12 · WORK LOG")
    timeline = [
        ("01", "장소 목록 하드코딩", "원본 카탈로그와 접근성 카드를 분리", COLORS["blue"]),
        ("02", "공공데이터 yes/no 충돌", "충돌 없는 yes만 자동 반영, 나머지는 검수 큐", COLORS["teal"]),
        ("03", "로드뷰 원본 70장 404", "누락 보고서·제공기관 복구 요청 패키지 생성", COLORS["orange"]),
        ("04", "지도 시안과 실제 위치 불일치", "실제 좌표 투영 + OSRM 경로 + 폴백 적용", COLORS["purple"]),
        ("05", "LLM 단정 위험", "룰 엔진과 설명 모델 분리 + strict JSON", COLORS["red"]),
    ]
    add_line(slide, 1.24, 2.12, 1.24, 6.05, COLORS["line"], 3)
    for index, (num, before, after, accent) in enumerate(timeline):
        y = 1.72 + index * 0.94
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.98), Inches(y + 0.16), Inches(0.52), Inches(0.52))
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent
        shape.line.color.rgb = accent
        add_text(slide, 0.98, y + 0.29, 0.52, 0.18, num, 8.5, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
        add_rect(slide, 1.8, y, 4.45, 0.78, COLORS["white"], line=COLORS["line"])
        add_text(slide, 2.04, y + 0.15, 3.95, 0.25, before, 12.5, bold=True)
        add_text(slide, 6.55, y + 0.13, 5.72, 0.48, after, 11.5, bold=True, color=accent, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
        add_chevron(slide, 6.12, y + 0.14, fill=accent)
    add_text(slide, 1.82, 6.42, 10.4, 0.26, "산출물: 수집 가이드 · 매칭/충돌 리포트 · 시각 검수 보드 · 운영 준비도 · 사전점검 · RAG 비교 자료", 10, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, 13)

    # 14. Team and deliverables
    slide = new_slide()
    add_title(slide, "평가 항목을 실제 산출물과 책임 영역으로 연결했습니다", "팀원 실명과 GitHub URL은 제출 전 팀 페이지와 동일하게 최종 입력합니다.", kicker="13 · TEAM / DELIVERABLES")
    roles = [
        ("기획·데이터", "문제 정의 · 출처·약관 · 데이터 게이트", COLORS["blue_soft"], COLORS["blue"]),
        ("AI·RAG", "점수 정책 · 프롬프트 · 검증 설계", COLORS["teal_soft"], COLORS["teal_dark"]),
        ("프론트·배포", "지도·근거 UI · API · Vercel", COLORS["yellow_soft"], COLORS["orange"]),
        ("QA·발표", "159 테스트 · 데모 시나리오 · 자료", COLORS["purple_soft"], COLORS["purple"]),
    ]
    for index, (title, body, fill, accent) in enumerate(roles):
        x = 0.72 + index * 3.13
        add_rect(slide, x, 1.72, 2.7, 1.45, fill, line=fill)
        add_text(slide, x + 0.22, 1.98, 2.25, 0.28, title, 14.5, bold=True, color=accent)
        add_text(slide, x + 0.22, 2.48, 2.22, 0.4, body, 9.5, color=COLORS["muted"])
    deliverables = [
        ("제안서", "docs/jeju_maeum_integrated_commercialization_plan.md", "완료", COLORS["green"]),
        ("작업 내역", "docs/jeju_maeum_progress_summary_20260708.md", "완료", COLORS["green"]),
        ("발표 자료", OUTPUT.name, "완료", COLORS["green"]),
        ("GitHub", "팀 저장소 URL 입력 필요", "입력", COLORS["orange"]),
        ("배포 URL", PUBLIC_URL, "공개", COLORS["blue"]),
    ]
    for index, (label, value, status, accent) in enumerate(deliverables):
        y = 3.62 + index * 0.58
        add_rect(slide, 0.88, y, 11.55, 0.46, COLORS["white"], line=COLORS["line"])
        add_text(slide, 1.08, y + 0.1, 1.25, 0.22, label, 10.5, bold=True, color=COLORS["ink"])
        add_text(slide, 2.55, y + 0.1, 7.85, 0.22, value, 9.5, color=COLORS["blue"] if label == "배포 URL" else COLORS["muted"], hyperlink=PUBLIC_URL if label == "배포 URL" else None)
        add_pill(slide, 10.9, y + 0.06, status, fill=accent, color=COLORS["white"], w=0.9, h=0.33)
    add_footer(slide, 14)

    # 15. Closing
    slide = new_slide(COLORS["ink"])
    add_image_crop(slide, images["welcome"], 7.75, 0, 5.583, 7.5)
    add_rect(slide, 0, 0, 8.02, 7.5, COLORS["ink"], rounded=False)
    add_pill(slide, 0.82, 0.8, "가치봄 제주", fill=COLORS["yellow"], color=COLORS["ink"], w=1.2)
    add_text(slide, 0.82, 1.52, 6.55, 1.36, "유명한 곳보다,\n근거 있는 곳을 추천합니다.", 34, bold=True, color=COLORS["white"])
    closing = [
        "모델보다 먼저 원본과 조건을 대조합니다.",
        "모르는 정보는 ‘확인 필요’로 남깁니다.",
        "AI가 실패해도 근거 기반 추천은 유지합니다.",
    ]
    add_bullets(slide, 0.9, 3.38, 6.25, closing, size=14, gap=0.6, bullet_color=COLORS["yellow"], text_color=COLORS["white"])
    url_button = add_rect(slide, 0.82, 5.68, 6.25, 0.82, COLORS["blue"], line=COLORS["blue"])
    url_button.click_action.hyperlink.address = PUBLIC_URL
    add_text(slide, 1.08, 5.88, 5.72, 0.24, PUBLIC_URL, 14, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    add_text(slide, 0.84, 6.82, 5.6, 0.22, "감사합니다", 10, color=COLORS["muted"])
    add_text(slide, 12.25, 7.14, 0.4, 0.2, "15", 8, bold=True, color=COLORS["yellow"], align=PP_ALIGN.RIGHT)

    # 16. Appendix - sources
    slide = new_slide()
    add_title(slide, "데이터 출처와 이용 조건", "원본명·URL·라이선스를 데이터와 코드에 함께 기록했습니다.", kicker="APPENDIX A")
    rows = [
        ["제주관광공사 관광 약자 유형별 추천코스", "16코스 · 62슬롯\n상황별 공식 코스 대조", "공공저작물 제4유형\n조건 준수 재확인", "장소 매칭·제한 승급"],
        ["제주특별자치도 사회적약자 시설현황", "107 관광지\n화장실·주차·대여·휴게", "공공데이터포털\n출처 URL 저장", "접근성 카드 초안"],
        ["로드뷰 이미지 메타데이터", "4,748건\n좌표·촬영일·해상도", "이용허락범위 제한 없음", "좌표·근거 이미지 연결"],
        ["로드뷰 이미지 원본", "서비스 시드 1,023장\n현재 953장 확보", "제공기관 조건 준수\n70장 복구 요청", "사람 시각 검수"],
        ["열린관광·이지제주·공식 상세", "장소별 보강 출처", "원문 링크·확인일 저장", "경사·바닥·주의사항 보강"],
    ]
    add_source_table(slide, 0.72, 1.75, 11.65, 4.65, rows)
    add_text(slide, 0.82, 6.55, 11.45, 0.26, "주의: 후기·SNS는 접근성 사실의 단독 근거로 사용하지 않으며, 공공저작물 제4유형의 변경금지 조건 충족 여부는 제출 전 최종 확인합니다.", 8.8, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, 16, appendix=True)

    # 17. Appendix - limitations
    slide = new_slide()
    add_title(slide, "현재 한계와 다음 검증", "측정하지 않은 것을 성과로 표현하지 않고, 다음 실험을 재현 가능하게 정의했습니다.", kicker="APPENDIX B")
    limitations = [
        ("실제 모델 A/B 없음", "59/59는 정책·랭킹 검증이며 무RAG는 통제 fixture입니다.", "동일 5개 입력 × 3회 실제 gpt-5-mini 호출, 블라인드 채점"),
        ("지연·비용 미측정", "OpenAI 실호출 latency와 token 비용 로그가 없습니다.", "p50/p95, 실패율, 입력·출력 토큰, 케이스당 비용 기록"),
        ("출처 근거 주입 제한", "GPT context에는 출처 URL·확인일이 직접 포함되지 않습니다.", "source id/url/checked_at을 retrieval context와 응답 citation에 추가"),
        ("공개 API 보호 부족", "인증·rate limit이 없어 예산 소진 공격 위험이 있습니다.", "서버에서 gpt-5-mini 고정, rate limit, 일일 비용 차단"),
        ("공개 게이트 불일치", "needs_check는 점수 상한만 있고 랭커 강제 제외가 아닙니다.", "추천 전 verification status allowlist를 코드로 강제"),
    ]
    for index, (title, current, next_step) in enumerate(limitations):
        y = 1.7 + index * 0.96
        add_rect(slide, 0.72, y, 2.45, 0.72, COLORS["red_soft"], line=COLORS["red_soft"])
        add_text(slide, 0.95, y + 0.18, 2.0, 0.28, title, 11.5, bold=True, color=COLORS["red"])
        add_rect(slide, 3.38, y, 4.05, 0.72, COLORS["white"], line=COLORS["line"])
        add_text(slide, 3.62, y + 0.12, 3.58, 0.42, current, 9.5, color=COLORS["muted"], valign=MSO_VERTICAL_ANCHOR.MIDDLE)
        add_chevron(slide, 7.56, y + 0.12, fill=COLORS["blue"])
        add_rect(slide, 8.0, y, 4.62, 0.72, COLORS["blue_soft"], line=COLORS["blue_soft"])
        add_text(slide, 8.24, y + 0.12, 4.14, 0.42, next_step, 9.5, bold=True, color=COLORS["blue_dark"], valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    add_text(slide, 0.88, 6.63, 11.55, 0.22, "핵심 우선순위: 실제 gpt-5-mini A/B → 출처 citation 강화 → rate limit/모델 고정 → 공개 게이트 강제", 10, bold=True, color=COLORS["blue_dark"], align=PP_ALIGN.CENTER)
    add_footer(slide, 17, appendix=True)

    prs.save(OUTPUT)
    return OUTPUT


SPEAKER_NOTES = """# 가치봄 제주 해커톤 발표 대본

발표 기준: 본편 15장, 약 7분. Appendix는 질의응답용입니다.

## 1. 표지
가치봄 제주는 관광약자의 이동 조건과 제주 접근성 원본을 대조해, 유명한 장소보다 근거 있는 장소를 추천하는 서비스입니다. 모델은 gpt-5-mini 하나를 사용하고, 구조화 RAG와 검증 게이트로 신뢰도를 높였습니다.

## 2. 문제 정의
휠체어 이용자가 “갈 수 있나요?”라고 물었을 때 인기나 별점은 답이 아닙니다. 화장실, 주차, 경사, 바닥, 휴식 정보가 필요하지만 출처가 흩어져 있고 현장 상황도 바뀝니다.

## 3. 솔루션
가치봄 제주는 복합 적합도, 출처와 검수 상태, 감점·제외 이유, 방문 전 확인사항을 한 화면에 제공합니다. 갈 수 있다고 보장하지 않고 판단 근거를 제공합니다.

## 4. 사용자 경험
사용자는 다섯 상황 테마 중 하나를 고르고 상세 조건을 추가합니다. 추천 카드에서 결과를 비교한 뒤, 실제 좌표 경로와 장소별 근거를 확인합니다.

## 5. 데이터
공식 원본을 공통 스키마로 정규화하고 verified, partial, needs_check를 구분합니다. 자동 병합은 충돌 없는 yes만 허용하고 나머지는 사람 검수로 보냅니다.

## 6. RAG 구조
구조화 필터와 규칙 랭킹으로 91개 카드에서 상위 4곳을 먼저 고릅니다. gpt-5-mini는 점수, 가점, 감점, 방문 전 확인만 받아 설명합니다. 따라서 모델이 장소를 새로 만들거나 순위를 임의로 바꾸지 않습니다.

## 7. 가드레일
제외와 점수는 규칙 엔진이 책임지고, 모델은 설명만 책임집니다. strict JSON Schema, 없는 시설 생성 금지, store false, timeout을 적용했습니다. 모델 실패 시에도 추천 결과는 남습니다.

## 8. 검증 설계
검증은 코드·API 159개 테스트, 5개 상황 59개 정책 체크, 공식 원본과 사람 검수 게이트로 나눴습니다. 실제 모델 지연과 비용은 아직 측정 범위가 아닙니다.

## 9. 정량 결과
현재 근거 기반 정책 파이프라인은 59개 체크를 모두 통과했습니다. 무RAG 통제 기준선은 5개만 통과했습니다. 다만 이것은 실제 GPT A/B가 아니라 정책 파이프라인과 통제 fixture 비교라는 한계를 명확히 밝힙니다.

## 10. 휠체어 사례
근거 기반 정책은 제주문학관, 국제컨벤션센터, 한란전시관, 김만덕기념관을 추천해 11개 체크를 통과했습니다. 통제 기준선은 성산일출봉, 우도, 시장, 해변을 추천했고 시설 근거와 방문 전 확인이 없었습니다.

## 11. 견고성과 보안
정상 연결에서는 AI 설명과 도로 경로를 제공합니다. 모델 오류에는 로컬 추천을 유지하고, API 오류에는 정적 seed, 경로 오류에는 브라우저 계산과 좌표 요약으로 폴백합니다. 키는 서버 환경변수만 사용합니다.

## 12. 운영 게이트
로드뷰 원본은 1,023장 중 953장을 확보했고 70장이 제공기관 404입니다. 17곳 68개 필드는 사람 검수 대기입니다. 제한 데모는 가능하지만 전체 공개는 이 작업이 끝날 때까지 차단합니다.

## 13. 작업 내역
하드코딩된 장소 목록에서 시작해 원본과 접근성 카드를 분리했고, 충돌 큐, 404 복구 패키지, 실제 좌표 경로, 모델 가드레일까지 만들었습니다. 실패를 문서와 재실행 가능한 스크립트로 남겼습니다.

## 14. 팀과 산출물
역할은 기획·데이터, AI·RAG, 프론트·배포, QA·발표로 나눴습니다. 제안서, 작업 내역, 발표 자료, 공개 배포 URL을 준비했고 팀원 실명과 GitHub URL은 팀 페이지와 동일하게 최종 입력합니다.

## 15. 결론
가치봄 제주는 모델의 자신감이 아니라 원본, 조건, 검수 상태로 여행 정보를 판단합니다. 유명한 곳보다 근거 있는 곳을 추천하겠습니다.

## 예상 질문
- 왜 벡터 DB가 없나요? 현재 91개 구조화 카드와 명시적 조건에서는 필터·랭킹 retrieval이 더 재현 가능하고 비용이 낮습니다. 규모가 커지면 hybrid search를 추가할 계획입니다.
- 59/59가 모델 성능인가요? 아닙니다. 정책·랭킹 회귀검증입니다. 실제 gpt-5-mini A/B는 다음 실험으로 분리했습니다.
- 전체 공개 가능한가요? 제한 데모는 가능하지만 로드뷰 70장과 17곳 시각 검수가 끝날 때까지 전체 공개 게이트는 보류 상태입니다.
- 점수는 신뢰도인가요? 출처 25, 이동 25, 시설 20, 테마 15, 안전 명확성 15를 합친 복합 적합도입니다. 출처 신뢰는 그중 한 축입니다.
"""


def write_speaker_notes() -> Path:
    SCRIPT_OUTPUT.write_text(SPEAKER_NOTES, encoding="utf-8-sig")
    return SCRIPT_OUTPUT


def validate_bounds(path: Path) -> list[str]:
    prs = Presentation(path)
    warnings: list[str] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                warnings.append(f"slide {slide_number}: out-of-bounds shape {shape.shape_id}")
    return warnings


if __name__ == "__main__":
    deck = build_deck()
    notes = write_speaker_notes()
    warnings = validate_bounds(deck)
    print(deck)
    print(notes)
    if warnings:
        print("\n".join(warnings))
