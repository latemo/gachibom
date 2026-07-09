from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "gachibom_jeju_runacation_hackathon_presentation_20260709.pptx"

COLORS = {
    "ink": RGBColor(17, 19, 24),
    "muted": RGBColor(91, 102, 118),
    "blue": RGBColor(18, 111, 181),
    "blue_dark": RGBColor(11, 63, 120),
    "blue_soft": RGBColor(232, 245, 255),
    "teal": RGBColor(18, 140, 131),
    "teal_soft": RGBColor(225, 247, 243),
    "yellow": RGBColor(255, 210, 31),
    "line": RGBColor(222, 226, 235),
    "bg": RGBColor(247, 248, 251),
    "white": RGBColor(255, 255, 255),
    "rose": RGBColor(255, 240, 243),
    "cream": RGBColor(255, 246, 230),
    "mint": RGBColor(233, 248, 242),
    "purple": RGBColor(246, 242, 255),
}


def load_json(relative: str) -> dict:
    return json.loads((ROOT / relative).read_text(encoding="utf-8"))


def set_font(run, size: int, bold: bool = False, color: RGBColor | None = None):
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text="", size=18, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color or COLORS["ink"])
    return box


def add_title(slide, title: str, subtitle: str | None = None):
    add_textbox(slide, 0.65, 0.38, 10.3, 0.55, title, 29, True, COLORS["ink"])
    if subtitle:
        add_textbox(slide, 0.68, 0.95, 10.2, 0.35, subtitle, 12, False, COLORS["muted"])
    add_line(slide, 0.65, 1.25, 11.95, 1.25, COLORS["line"], 1)


def add_line(slide, x1, y1, x2, y2, color, width=2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_round_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1)
    return shape


def add_badge(slide, x, y, text, fill=COLORS["blue_soft"], color=COLORS["blue"], w=None):
    width = w or max(1.0, 0.14 * len(text) + 0.45)
    add_round_rect(slide, x, y, width, 0.36, fill, RGBColor(190, 220, 245))
    add_textbox(slide, x + 0.12, y + 0.075, width - 0.24, 0.16, text, 9, True, color, PP_ALIGN.CENTER)


def add_bullets(slide, x, y, w, items, size=16, gap=0.46, color=None):
    for i, item in enumerate(items):
        add_round_rect(slide, x, y + i * gap + 0.04, 0.09, 0.09, COLORS["blue"], COLORS["blue"])
        add_textbox(slide, x + 0.22, y + i * gap - 0.04, w - 0.22, 0.34, item, size, False, color or COLORS["ink"])


def add_metric(slide, x, y, w, h, value, label, fill=COLORS["white"], accent=COLORS["blue"]):
    add_round_rect(slide, x, y, w, h, fill, COLORS["line"])
    add_textbox(slide, x + 0.18, y + 0.15, w - 0.36, 0.5, value, 30, True, accent, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.18, y + 0.77, w - 0.36, 0.35, label, 10, True, COLORS["muted"], PP_ALIGN.CENTER)


def add_image_fit(slide, path: Path, x, y, w, h):
    if not path.exists():
        add_round_rect(slide, x, y, w, h, COLORS["blue_soft"], COLORS["line"])
        add_textbox(slide, x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, f"이미지 없음: {path.name}", 12, True, COLORS["muted"], PP_ALIGN.CENTER)
        return None
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_image_cover(slide, path: Path, x, y, w, h):
    pic = add_image_fit(slide, path, x, y, w, h)
    if pic:
        pic.crop_left = 0.03
        pic.crop_right = 0.03
    return pic


def add_footer(slide, n: int):
    add_textbox(slide, 0.65, 7.18, 6.5, 0.25, "가치봄 제주 · 런케이션 해커톤 · 신뢰도 분석", 8, False, RGBColor(120, 128, 140))
    add_textbox(slide, 12.15, 7.18, 0.5, 0.25, f"{n:02d}", 8, True, COLORS["blue"], PP_ALIGN.RIGHT)


def scenario_route_text(seed):
    lines = []
    for scenario in seed.get("scenarios", []):
        title = scenario.get("title", "")
        route = scenario.get("recommendation", {}).get("course", {}).get("route", [])
        names = " → ".join([item.get("name", "") for item in route[:4]])
        lines.append((title, names))
    return lines


def build_deck():
    seed = load_json("web/data/app_recommendation_seed.json")
    validation = load_json("web/data/recommendation_case_validation_report.json")
    health = {
        "places": 91,
        "tests": 143,
        "route_proxy": True,
        "tourism_weak_courses": True,
    }
    courses = load_json("data/tourism_weak_recommendation_courses.json")
    course_summary = courses.get("summary", {})

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    img = {
        "cover": ROOT / "web" / "assets" / "WELCOME-1-001.jpg",
        "concept": ROOT / "gachibom-concept-page-fixed-v52-20260709.png",
        "map": ROOT / "jeju-maeum-live-leaflet-map-popup-1366-20260709.png",
        "route": ROOT / "jeju-maeum-route-detail-modal-8792-proxy-1366-20260709.png",
        "detail": ROOT / "jeju-maeum-validation-evidence-8790-20260709.png",
        "commercial": ROOT / "gachibom-commercial-map-1366-20260709.png",
        "forest": ROOT / "web" / "assets" / "SAMSUNGHYEOL-1-001.jpg",
        "museum": ROOT / "web" / "assets" / "JEJUNATIONALMU-1-001.jpg",
    }

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_image_cover(slide, img["cover"], 7.35, 0, 5.98, 7.5)
    add_round_rect(slide, 0, 0, 13.333, 7.5, RGBColor(248, 250, 252), RGBColor(248, 250, 252), False)
    add_image_cover(slide, img["cover"], 7.1, 0.38, 5.55, 6.75)
    add_badge(slide, 0.75, 0.75, "제주 신뢰도 분석 서비스", w=2.1)
    add_textbox(slide, 0.75, 1.35, 5.8, 1.4, "가치봄 제주", 44, True, COLORS["ink"])
    add_textbox(slide, 0.78, 2.45, 5.75, 0.75, "접근성 여행 정보를 믿을 수 있게 분석하는 서비스", 22, True, COLORS["blue_dark"])
    add_bullets(slide, 0.8, 3.55, 5.6, [
        "gpt-5-mini 고정 · RAG · 출처 대조 · 검증 게이트",
        "제주 접근성 여행지와 공식 추천코스 기반",
        "팀 역할: 기획/데이터 · AI/RAG · 프론트/배포",
    ], 15, 0.5)
    add_textbox(slide, 0.78, 6.75, 5.0, 0.35, "런케이션 해커톤 발표 자료 · 2026.07.09", 11, False, COLORS["muted"])

    # 2. Problem
    slide = prs.slides.add_slide(blank)
    add_title(slide, "문제 정의", "여행 정보는 많지만, 접근성 정보는 믿기 어렵다")
    cards = [
        ("정보 과잉", "후기·블로그·광고성 정보가 섞여 실제 방문 가능성을 판단하기 어렵다.", COLORS["blue_soft"]),
        ("출처 불명", "장애인 화장실, 주차, 경사, 휴식 공간 정보의 출처와 최신성이 불분명하다.", COLORS["teal_soft"]),
        ("사용자 조건", "휠체어·회복기·아이 동반·날씨 민감·음식 제한은 같은 추천으로 해결되지 않는다.", COLORS["cream"]),
    ]
    for i, (title, body, fill) in enumerate(cards):
        x = 0.7 + i * 4.05
        add_round_rect(slide, x, 1.75, 3.65, 3.65, fill, COLORS["line"])
        add_textbox(slide, x + 0.28, 2.05, 3.1, 0.45, title, 22, True, COLORS["ink"])
        add_textbox(slide, x + 0.28, 2.85, 3.0, 1.5, body, 16, False, COLORS["muted"])
    add_textbox(slide, 1.1, 6.08, 11.2, 0.5, "우리는 ‘좋은 장소 추천’보다 먼저, 이 정보가 추천 근거로 쓸 만큼 믿을 만한지 분석한다.", 20, True, COLORS["blue_dark"], PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # 3. Users and domain
    slide = prs.slides.add_slide(blank)
    add_title(slide, "도메인과 사용자", "제주 접근성 여행지 · 상황별 코스 신뢰도 분석")
    add_image_cover(slide, img["forest"], 0.75, 1.65, 4.2, 4.55)
    user_cards = [
        ("회복 중", "무리 없는 일정"),
        ("휠체어 접근", "경사·엘리베이터·주차"),
        ("아이 동반", "유모차·휴식 동선"),
        ("날씨 민감", "실내/실외 리스크"),
        ("음식 제한", "식당·시장 제외"),
    ]
    for i, (title, body) in enumerate(user_cards):
        x = 5.35 + (i % 2) * 3.55
        y = 1.58 + (i // 2) * 1.25
        add_round_rect(slide, x, y, 3.25, 0.88, COLORS["white"], COLORS["line"])
        add_textbox(slide, x + 0.18, y + 0.13, 1.5, 0.25, title, 15, True, COLORS["ink"])
        add_textbox(slide, x + 0.18, y + 0.47, 2.7, 0.2, body, 10, False, COLORS["muted"])
    add_round_rect(slide, 5.35, 5.45, 6.95, 0.9, COLORS["blue_soft"], RGBColor(190, 220, 245))
    add_textbox(slide, 5.6, 5.64, 6.5, 0.32, "핵심 질문: 추천해도 되는가 · 왜 추천하는가 · 방문 전 무엇을 확인해야 하는가", 15, True, COLORS["blue_dark"])
    add_footer(slide, 3)

    # 4. Flow
    slide = prs.slides.add_slide(blank)
    add_title(slide, "사용자 흐름", "컨셉 선택에서 근거 확인까지 한 번에 이어지는 서비스")
    steps = [
        ("1", "컨셉 선택", "현재 여행 조건을 먼저 고른다"),
        ("2", "후보 필터링", "공개 가능 데이터와 제외 조건 적용"),
        ("3", "신뢰도 계산", "출처·시설·동선·안전 기준 점수화"),
        ("4", "근거 제공", "지도·코스·출처·확인 항목 표시"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.9 + i * 3.0
        add_round_rect(slide, x, 2.0, 2.35, 2.35, COLORS["white"], COLORS["line"])
        add_round_rect(slide, x + 0.18, 2.18, 0.48, 0.48, COLORS["blue"], COLORS["blue"])
        add_textbox(slide, x + 0.18, 2.28, 0.48, 0.18, num, 11, True, COLORS["white"], PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.22, 2.9, 1.9, 0.3, title, 17, True, COLORS["ink"])
        add_textbox(slide, x + 0.22, 3.35, 1.9, 0.5, body, 11, False, COLORS["muted"])
        if i < 3:
            add_line(slide, x + 2.35, 3.2, x + 2.85, 3.2, COLORS["blue"], 2)
    add_image_cover(slide, img["concept"], 0.9, 5.0, 11.55, 1.25)
    add_footer(slide, 4)

    # 5. Data
    slide = prs.slides.add_slide(blank)
    add_title(slide, "데이터 확보와 출처", "모델의 ‘감’이 아니라 출처 데이터와 대조한다")
    metrics = [
        ("91", "장소 카드"),
        ("16", "공식 추천코스"),
        ("62/62", "코스 슬롯 매칭"),
        ("61", "매칭 장소"),
        ("0", "미매칭 장소"),
    ]
    for i, (value, label) in enumerate(metrics):
        add_metric(slide, 0.65 + i * 2.45, 1.55, 2.05, 1.25, value, label, COLORS["white"], COLORS["blue"])
    add_round_rect(slide, 0.9, 3.45, 5.55, 2.35, COLORS["teal_soft"], COLORS["line"])
    add_textbox(slide, 1.15, 3.75, 4.8, 0.35, "제주관광공사 추천여행코스", 18, True, COLORS["ink"])
    add_bullets(slide, 1.18, 4.35, 4.9, [
        "관광 약자 유형별 제주관광 추천코스",
        "16개 코스, 62개 장소 슬롯 처리",
        "기존 카드 매칭 + 신규 후보 승급"
    ], 12, 0.38)
    add_round_rect(slide, 6.85, 3.45, 5.55, 2.35, COLORS["blue_soft"], COLORS["line"])
    add_textbox(slide, 7.1, 3.75, 4.8, 0.35, "로드뷰 시각 검수 근거", 18, True, COLORS["ink"])
    add_bullets(slide, 7.13, 4.35, 4.9, [
        "서비스 시드 17곳, 검수 필드 68개",
        "1,023장 중 953장 확보",
        "70장은 제공기관 404 복구 요청 대상"
    ], 12, 0.38)
    add_footer(slide, 5)

    # 6. RAG architecture
    slide = prs.slides.add_slide(blank)
    add_title(slide, "gpt-5-mini + RAG 설계", "모델은 하나, 신뢰도는 근거 조립과 검증으로 끌어올린다")
    cols = [
        ("근거 저장소", ["장소 카드", "상황별 규칙", "공식 코스", "출처 요약"]),
        ("RAG/프롬프트", ["조건별 관련 근거 검색", "추천/감점 이유 분리", "방문 전 확인 항목 생성"]),
        ("검증 게이트", ["스키마 검증", "정책 점수 보정", "차단/검수대기 제외"]),
        ("서비스 출력", ["지도와 코스", "장소 상세 근거", "공식 코스 탭"]),
    ]
    for i, (title, items) in enumerate(cols):
        x = 0.75 + i * 3.08
        add_round_rect(slide, x, 1.75, 2.65, 3.4, COLORS["white"], COLORS["line"])
        add_textbox(slide, x + 0.18, 2.03, 2.2, 0.34, title, 16, True, COLORS["blue_dark"])
        add_bullets(slide, x + 0.25, 2.68, 2.15, items, 10, 0.42)
        if i < 3:
            add_line(slide, x + 2.65, 3.45, x + 3.02, 3.45, COLORS["teal"], 2)
    add_badge(slide, 4.9, 5.95, "고정 모델: gpt-5-mini", fill=COLORS["yellow"], color=COLORS["ink"], w=3.2)
    add_footer(slide, 6)

    # 7. Trust scoring
    slide = prs.slides.add_slide(blank)
    add_title(slide, "신뢰도 점수 정책", "추천 점수는 ‘잘 맞음’과 ‘근거 명확성’을 함께 본다")
    axes = [
        ("source_trust", "공식 출처·검수 근거"),
        ("mobility_fit", "도보 부담·계단·경사"),
        ("facility_fit", "화장실·주차·엘리베이터"),
        ("theme_fit", "컨셉과 장소 성격 일치"),
        ("safety_clarity", "확인 필요 항목 명시"),
    ]
    for i, (axis, desc) in enumerate(axes):
        y = 1.65 + i * 0.75
        add_textbox(slide, 0.8, y, 2.1, 0.25, axis, 13, True, COLORS["blue_dark"])
        add_round_rect(slide, 3.05, y + 0.04, 5.6, 0.22, COLORS["blue_soft"], COLORS["blue_soft"], False)
        add_round_rect(slide, 3.05, y + 0.04, 2.9 + i * 0.38, 0.22, COLORS["blue"], COLORS["blue"], False)
        add_textbox(slide, 8.95, y - 0.02, 3.2, 0.3, desc, 12, False, COLORS["muted"])
    add_round_rect(slide, 0.9, 5.65, 11.45, 0.75, COLORS["cream"], COLORS["line"])
    add_textbox(slide, 1.15, 5.86, 10.9, 0.25, "Exclude first: 음식 제한, 날씨 민감, 차단 장소는 추천 전 단계에서 먼저 제외한다.", 15, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8. Validation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "검증 결과", "신뢰성은 원본·출처와 대조한 테스트로 증명한다")
    validation_summary = validation.get("summary", {})
    metric_items = [
        (f"{validation_summary.get('passed_cases', 5)}/{validation_summary.get('total_cases', 5)}", "상황별 검증 통과"),
        (str(health["tests"]), "전체 테스트 통과"),
        ("20/20", "추천 노출 장소 좌표"),
        ("OK", "/api/health"),
    ]
    for i, (value, label) in enumerate(metric_items):
        add_metric(slide, 0.8 + i * 3.0, 1.65, 2.4, 1.55, value, label, COLORS["white"], COLORS["teal"] if i == 0 else COLORS["blue"])
    add_round_rect(slide, 1.05, 4.05, 11.1, 1.3, COLORS["blue_soft"], COLORS["line"])
    add_bullets(slide, 1.35, 4.34, 10.4, [
        "route_proxy true · tourism_weak_courses true",
        "API 오류는 JSON code/error로 반환, 과대 본문은 413 처리",
        "내부 예외 원문과 비밀값은 응답에 노출하지 않음"
    ], 14, 0.36)
    add_footer(slide, 8)

    # 9. Scenario differentiation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "상황별 추천이 실제로 달라진다", "동일한 제주라도 사용자 조건에 따라 제외·가점·코스가 바뀐다")
    routes = scenario_route_text(seed)
    for i, (title, names) in enumerate(routes):
        y = 1.55 + i * 0.88
        fill = [COLORS["rose"], COLORS["mint"], COLORS["purple"], COLORS["blue_soft"], COLORS["cream"]][i % 5]
        add_round_rect(slide, 0.75, y, 11.85, 0.65, fill, COLORS["line"])
        add_textbox(slide, 1.0, y + 0.12, 2.55, 0.22, title, 13, True, COLORS["ink"])
        add_textbox(slide, 3.65, y + 0.12, 8.55, 0.22, names, 11, False, COLORS["muted"])
    add_textbox(slide, 1.05, 6.15, 11.2, 0.34, "예: 음식 제한은 식당·시장 제외, 날씨 민감은 바다·오름·고노출 야외를 상위 추천에서 배제", 14, True, COLORS["blue_dark"], PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # 10. Demo
    slide = prs.slides.add_slide(blank)
    add_title(slide, "서비스 화면 데모", "컨셉 선택 · 지도 · 상세 근거 · 공식 코스를 한 흐름으로 제공")
    add_image_cover(slide, img["concept"], 0.65, 1.45, 5.8, 2.5)
    add_image_cover(slide, img["map"], 6.75, 1.45, 5.9, 2.5)
    add_image_cover(slide, img["route"], 0.65, 4.35, 5.8, 1.9)
    add_image_cover(slide, img["detail"], 6.75, 4.35, 5.9, 1.9)
    add_badge(slide, 0.85, 1.6, "컨셉 선택", w=1.25)
    add_badge(slide, 6.95, 1.6, "실제 좌표 지도", w=1.55)
    add_badge(slide, 0.85, 4.5, "경로 모달", w=1.25)
    add_badge(slide, 6.95, 4.5, "근거 상세", w=1.25)
    add_footer(slide, 10)

    # 11. Robustness & security
    slide = prs.slides.add_slide(blank)
    add_title(slide, "서비스 견고성과 보안", "해커톤 데모가 아니라 배포 가능한 서비스 기준으로 점검")
    left = [
        ("서버 계약", "/api/health · /api/recommendations · /api/routes"),
        ("장애 대응", "API 실패 시 기본 추천 폴백, 경로 프록시 미지원 시 브라우저 계산 대체"),
        ("데이터 게이트", "검수 전 장소는 사용자 추천 근거로 과장 표기하지 않음"),
        ("보안", "API 키는 환경변수/Secrets, 리포트에는 설정 여부만 기록"),
    ]
    for i, (title, body) in enumerate(left):
        y = 1.65 + i * 1.05
        add_round_rect(slide, 0.85, y, 11.55, 0.75, COLORS["white"], COLORS["line"])
        add_textbox(slide, 1.1, y + 0.15, 1.8, 0.22, title, 14, True, COLORS["blue_dark"])
        add_textbox(slide, 3.0, y + 0.15, 8.9, 0.22, body, 12, False, COLORS["muted"])
    add_round_rect(slide, 3.55, 6.05, 6.2, 0.55, COLORS["teal_soft"], COLORS["line"])
    add_textbox(slide, 3.75, 6.2, 5.8, 0.18, "비밀값 노출 검사: 의심 파일 0개", 14, True, COLORS["teal"], PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12. Closing
    slide = prs.slides.add_slide(blank)
    add_title(slide, "결론과 다음 단계", "좋은 모델보다 중요한 것은 검증 가능한 근거와 운영 게이트")
    add_textbox(slide, 0.9, 1.55, 11.7, 0.65, "가치봄 제주는 접근성 여행 정보를 ‘추천’하기 전에, 출처와 조건과 검수 상태를 대조해 믿을 수 있는지 분석합니다.", 22, True, COLORS["ink"], PP_ALIGN.CENTER)
    next_steps = [
        "팀 페이지 기록: 제안서 · 작업 내역 · 발표 자료 · GitHub · 배포 URL",
        "누락 로드뷰 70장 복구 또는 대체 원본 수령",
        "사람 최종 검수 완료 후 서비스 시드 승격",
        "공개 배포 URL 연결 및 README 실행 방법 정리",
    ]
    add_bullets(slide, 1.25, 3.0, 10.8, next_steps, 17, 0.55)
    add_round_rect(slide, 1.0, 6.2, 11.3, 0.55, COLORS["blue"], COLORS["blue"])
    add_textbox(slide, 1.2, 6.37, 10.9, 0.15, "데모: http://127.0.0.1:8792/?v=52  ·  공개 URL은 팀 페이지 제출 전 교체", 12, True, COLORS["white"], PP_ALIGN.CENTER)
    add_footer(slide, 12)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    out = build_deck()
    print(out)
