import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]


class RagComparisonPresentationTests(unittest.TestCase):
    def test_cli_builds_presentation_with_key_summary(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "rag_comparison.pptx"
            result = subprocess.run(
                [
                    sys.executable,
                    str(ROOT / "scripts" / "build_rag_comparison_presentation.py"),
                    "--comparison-json",
                    str(ROOT / "data" / "rag_comparison_report.json"),
                    "--no-rag-json",
                    str(ROOT / "data" / "no_rag_baseline_validation_report.json"),
                    "--output",
                    str(output),
                ],
                cwd=ROOT,
                capture_output=True,
                text=True,
                check=False,
            )

            self.assertEqual(result.returncode, 0, result.stderr)
            self.assertTrue(output.exists())
            presentation = Presentation(output)
            self.assertEqual(len(presentation.slides), 8)
            text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            self.assertIn("RAG 사용/미사용 비교", text)
            self.assertIn("59/59", text)
            self.assertIn("5/59", text)
            self.assertIn("무RAG 실패 패턴", text)


if __name__ == "__main__":
    unittest.main()
